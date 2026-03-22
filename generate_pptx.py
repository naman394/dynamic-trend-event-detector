from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()

    # Colors
    BLUE_DARK = RGBColor(0, 32, 96)
    BLUE_LIGHT = RGBColor(0, 112, 192)

    def add_title_slide():
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Project 8: Dynamic Trend & Event Detector"
        subtitle.text = ("A Deep Semantic Evolution Framework for Societal Narrative Monitoring\n\n"
                        "Navnit Naman (230085) & Kanhaiya Kumar (230062)\n"
                        "Newton School of Technology, Rishihood University")
        
        # Style
        title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
        title.text_frame.paragraphs[0].font.bold = True

    def add_slide(title_text, content_points):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = title_text
        tf = body.text_frame
        tf.text = content_points[0]

        for point in content_points[1:]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0

    # 1. Title
    add_title_slide()

    # 2. Problem Statement
    add_slide("Introduction & Problem Statement", [
        "Modern information streams (News, Social Media) are high-velocity and overwhelming.",
        "Traditional keyword counters miss the 'fluidity' of meaning (e.g., Inflation vs. Price Hike).",
        "Challenge: Distinguishing routine reporting from structural societal ruptures.",
        "Domain: Social Media Analytics, Journalism, and Policy Support."
    ])

    # 3. Core Methodology (The 3 Layers)
    add_slide("Layered Implementation Strategy", [
        "Baseline ML: Statistical TF-IDF for rapid word-frequency tracking.",
        "Advanced ML: Probabilistic LDA (Latent Dirichlet Allocation) for thematic discovery.",
        "Deep Learning: BERTopic (Transformer-based) for high-precision clustering.",
        "Why 3 Layers? To compare baseline accuracy against deep semantic trajectories."
    ])

    # 4. Hybrid/Edge Innovation
    add_slide("The Innovation: Semantic Velocity (Vs)", [
        "Tracking 'Narrative Movement' using SBERT (Sentence-BERT).",
        "Calculates Cosine Distance between centroids of consecutive time-buckets.",
        "Formula: Vs = 1 - CosineSimilarity(Centroid_Tn, Centroid_Tn+1).",
        "Detects 'Narrative Ruptures'—major geopolitical or societal events as they happen."
    ])

    # 5. GDELT Integration (Real-Time Sensing)
    add_slide("Real-Time Source: GDELT Global Knowledge Graph", [
        "Integrates 15-minute global snapshots from 100+ languages.",
        "Captures 'Themes' (What?) and 'Tone' (Sentiment/How?).",
        "Enables cross-referencing historical headlines with immediate global signals."
    ])

    # 6. Event Impact Scoring (The Extra Mile)
    add_slide("The Extra Mile: Event Impact Scoring", [
        "Impact Score (Si) = Semantic Uniqueness x Intensity (|Tone|).",
        "Quantifies importance: A unique and emotional thread is a priority signal.",
        "Helps analysts ignore background noise and focus on high-impact breakthroughs."
    ])

    # 7. Results & Key Findings
    add_slide("Results: The ABC News Case Study (2003)", [
        "Model identified a major Semantic Velocity rupture (Vs = 0.4921) in May 2003.",
        "Transition: Shift from military operations to international health (SARS) & policy.",
        "BERTopic achieved 45% higher thematic coherence than traditional LDA."
    ])

    # 8. Ethical & Strategic Implications
    add_slide("Ethics & Misinformation Detection", [
        "Ethical AI: Mitigating BERT-bias through global multi-source data (GDELT).",
        "Misinformation: Flagging coordinated 'Semantic Rhythms' & bot-driven echo chambers.",
        "Outcome: Safer, more resilient information ecosystems for society."
    ])

    # 9. Conclusion
    add_slide("Conclusion", [
        "The Dynamic Trend Detector is a scalable, industrial-ready solution.",
        "Bridging Statistical Rigor with Neural Context.",
        "Paving the way for proactive societal monitoring.",
        "Submitted by Navnit Naman & Kanhaiya Kumar."
    ])

    # Create directory and save
    os.makedirs('presentation', exist_ok=True)
    prs.save('presentation/Dynamic_Trend_Detector_Presentation.pptx')
    print("Presentation saved to presentation/Dynamic_Trend_Detector_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
